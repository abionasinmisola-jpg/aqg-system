import fitz  # pymupdf
from docx import Document
from pptx import Presentation
import openpyxl
import os
import re


def extract_text(file_path: str, file_type: str) -> str:
    """Extract text from any supported file type"""
    extractors = {
        "pdf":  extract_from_pdf,
        "docx": extract_from_docx,
        "txt":  extract_from_txt,
        "pptx": extract_from_pptx,
        "xlsx": extract_from_xlsx,
    }

    extractor = extractors.get(file_type.lower())
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")

    raw_text = extractor(file_path)
    return clean_text(raw_text)


def extract_from_pdf(file_path: str) -> str:
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text


def extract_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])


def extract_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_from_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text += row_text + "\n"
    return text


def clean_text(text: str) -> str:
    """Clean and normalize extracted text"""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove special characters but keep punctuation
    text = re.sub(r'[^\w\s.,!?;:()\-\'\"]+', ' ', text)
    # Remove very short lines
    lines = [line.strip() for line in text.split('.') if len(line.strip()) > 20]
    return '. '.join(lines)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """Split text into overlapping chunks for RAG"""
    words = text.split()
    chunks = []
    start = 0

    while start < len(words):
        end = start + chunk_size
        chunk = ' '.join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap

    return chunks